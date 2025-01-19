import os
import requests
import io
from PyPDF2 import PdfReader
from PIL import Image
import pytesseract
from moviepy import *
import tempfile
from pydub import AudioSegment
import whisper
import logging
from docx import Document
import openpyxl
from pptx import Presentation
from logging.handlers import RotatingFileHandler

# Load Whisper model
whisper_model = whisper.load_model("base")

# Configure rotating file handler for logging
log_file = "media_processing.log"
log_handler = RotatingFileHandler(
    log_file, maxBytes=5 * 1024 * 1024, backupCount=5  # 5 MB per file, keep 5 backups
)
log_handler.setLevel(logging.DEBUG)
log_formatter = logging.Formatter("%(asctime)s - %(levelname)s - %(message)s")
log_handler.setFormatter(log_formatter)

# Set up the logger
logger = logging.getLogger()
logger.setLevel(logging.DEBUG)
logger.addHandler(log_handler)


def convert_to_wav(file_path):
    """Convert audio to WAV format."""
    audio = AudioSegment.from_file(file_path)
    wav_path = file_path.replace(".mp3", ".wav").replace(".ogg", ".wav")
    audio.export(wav_path, format="wav")
    return wav_path

def pdf_to_text(file_bytes):
    response_body = ""
    try:
        reader = PdfReader(file_bytes)
        if reader.is_encrypted:
            logger.warning("PDF file is encrypted. Attempting to decrypt...")
            # Attempt decryption (if you have the password or expect no password)
            try:
                reader.decrypt("")
            except Exception as e:
                logger.error(f"Failed to decrypt PDF: {e}")
                return "Error: Unable to process encrypted PDF."
            
        # Extract text
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        logger.info("PDF processed successfully.")
        return text
    except Exception as e:
        logger.error(f"Failed to process PDF: {e}")
        return f"Error: {e}"

def image_to_text(file_bytes):
    response_body = ""
    try:
        image = Image.open(file_bytes)
        text = pytesseract.image_to_string(image)
        response_body += f"\nExtracted text from Image:\n{text}"
    except Exception as e:
        response_body += f"\nFailed to process Image: {str(e)}"

    return response_body

def audio_to_text(file_bytes):
    response_body = ""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as temp_audio_file:
            temp_audio_file.write(file_bytes.read())
            temp_audio_file.flush()
            wav_path = temp_audio_file.name

            # Transcribe using Whisper
            result = whisper_model.transcribe(wav_path)
            text = result.get("text", "").strip()
            response_body += f"\nExtracted text from Audio:\n{text}"
    except Exception as e:
        response_body += f"\nFailed to process Audio: {str(e)}"

    return response_body

def video_to_text(file_bytes,media_url=None):
    response_body = ""
    try:

        try:
            response = requests.get(
                media_url,
                auth=(os.getenv("TWILIO_ACCOUNT_SID"), os.getenv("TWILIO_AUTH_TOKEN")),
                timeout=15  # Set timeout
            )
            response.raise_for_status()

            file_size = len(response.content)
            print(f"File Size: {file_size} bytes")  # Debugging

            if file_size > 50 * 1024 * 1024:  # Limit to 50MB
                raise Exception("File is too large to process.")
            
            file_bytes = io.BytesIO(response.content)
            # Process the file as before...

        except requests.exceptions.RequestException as e:
            response_body += f"\nFailed to download media file: {e}"
            print(f"Download Error: {e}")

        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as temp_video_file:
            temp_video_file.write(file_bytes.read())
            temp_video_file.flush()
            video_path = temp_video_file.name

            # Extract audio from video
            with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as temp_audio_file:
                video_clip = VideoFileClip(video_path)
                video_clip.audio.write_audiofile(temp_audio_file.name)
                wav_path = temp_audio_file.name

                # Transcribe using Whisper
                result = whisper_model.transcribe(wav_path)
                text = result.get("text", "").strip()
                response_body += f"\nExtracted text from Video:\n{text}"
    except Exception as e:
        response_body += f"\nFailed to process Video: {str(e)}"

    return response_body

def word_to_text(file_bytes):
    try:
        doc = Document(file_bytes)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return f"Extracted text from Word Document:\n{text}"
    except Exception as e:
        logger.error(f"Failed to process Word file: {e}")
        return f"Error processing Word file: {e}"

def excel_to_text(file_bytes):
    try:
        workbook = openpyxl.load_workbook(file_bytes)
        text = ""
        for sheet in workbook.worksheets:
            text += f"\nSheet: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                text += " ".join(map(str, row)) + "\n"
        return f"Extracted text from Excel Document:\n{text}"
    except Exception as e:
        logger.error(f"Failed to process Excel file: {e}")
        return f"Error processing Excel file: {e}"

def ppt_to_text(file_bytes):
    """
    Extract text from a PowerPoint file.
    """
    try:
        logger.debug("Starting PowerPoint file processing...")
        presentation = Presentation(file_bytes)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n"

        logger.info("PowerPoint file processed successfully.")
        
        # Log the length of the extracted text
        logger.debug(f"Extracted text length: {len(text)} characters.")
        
        # Optionally log a snippet of the text (first 500 characters)
        logger.debug(f"Extracted text snippet: {text[:500]}")

        return f"Extracted text from PowerPoint Document:\n{text}"

    except Exception as e:
        logger.error(f"Failed to process PowerPoint file: {e}")
        print(f"PowerPoint Processing Error: {e}")  # Debug for terminal
        return f"Error processing PowerPoint file: {e}"